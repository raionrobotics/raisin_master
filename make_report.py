"""
RAISIN 단위 테스트 보고서 PPT 생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 정의 ──────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x3E)
BLUE   = RGBColor(0x1A, 0x5F, 0xA8)
GREEN  = RGBColor(0x1A, 0x8C, 0x5A)
RED    = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF4, 0xF6, 0xF8)
DGRAY  = RGBColor(0x5D, 0x6D, 0x7E)
BLACK  = RGBColor(0x17, 0x20, 0x2A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]   # 완전 빈 레이아웃

def add_rect(slide, l, t, w, h, fill, alpha=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s

def add_text_box(slide, text, l, t, w, h,
                 size=18, bold=False, color=BLACK,
                 align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.35, NAVY)
    add_text_box(slide, title, 0.4, 0.08, 12, 0.7,
                 size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.4, 0.85, 12, 0.4,
                     size=13, bold=False, color=RGBColor(0xAA, 0xC4, 0xFF),
                     align=PP_ALIGN.LEFT)

def footer(slide, page, total=5):
    add_rect(slide, 0, 7.2, 13.33, 0.3, NAVY)
    add_text_box(slide, "RAISIN Unit Test Report  |  2026-05-22",
                 0.3, 7.21, 10, 0.28, size=9, color=WHITE)
    add_text_box(slide, f"{page} / {total}",
                 12.5, 7.21, 0.8, 0.28, size=9, color=WHITE, align=PP_ALIGN.RIGHT)

def bullet_block(slide, items, l, t, w, h, title=None, title_color=BLUE, bullet="•"):
    if title:
        add_text_box(slide, title, l, t, w, 0.32,
                     size=13, bold=True, color=title_color)
        t += 0.34
        h -= 0.34
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(11.5)
        run.font.color.rgb = BLACK

def tag(slide, text, l, t, color=BLUE):
    add_rect(slide, l, t, len(text)*0.085+0.1, 0.25, color)
    add_text_box(slide, text, l+0.05, t+0.02, len(text)*0.085, 0.22,
                 size=9, bold=True, color=WHITE)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 슬라이드 1: 표지
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s1 = prs.slides.add_slide(blank)
add_rect(s1, 0, 0, 13.33, 7.5, NAVY)
add_rect(s1, 0, 3.8, 13.33, 3.7, BLUE)

add_text_box(s1, "RAISIN 패키지", 0.6, 0.7, 12, 0.9,
             size=22, bold=False, color=RGBColor(0xAA,0xC4,0xFF))
add_text_box(s1, "단위 테스트 구현 및 결과 보고서", 0.6, 1.5, 12, 1.1,
             size=38, bold=True, color=WHITE)

add_rect(s1, 0.6, 2.9, 2.0, 0.04, RGBColor(0x64,0xB5,0xF6))

add_text_box(s1, "테스트 대상 패키지", 0.6, 4.1, 6, 0.4,
             size=13, bold=True, color=WHITE)
pkgs = [
    ("raisin_autonomy_plugin", "map_graph 알고리즘 (A*, 투영, 경로 정제)"),
    ("raisin_ublox_plugin",    "Fletcher 체크섬 / NMEA 체크섬"),
]
for i, (pkg, desc) in enumerate(pkgs):
    yy = 4.6 + i*0.55
    add_rect(s1, 0.6, yy, 12.3, 0.42, RGBColor(0x0A, 0x3A, 0x7A))
    add_text_box(s1, pkg,  0.75, yy+0.03, 4.5, 0.35, size=11, bold=True,  color=RGBColor(0x64,0xB5,0xF6))
    add_text_box(s1, desc, 5.3,  yy+0.03, 7.3, 0.35, size=11, bold=False, color=WHITE)

add_text_box(s1, "빌드 환경: GTest / Ninja / C++20  |  빌드 타입: Release  |  작성: 2026-05-22",
             0.6, 6.8, 12, 0.35, size=10, color=RGBColor(0xAA,0xC4,0xFF))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 슬라이드 2: 테스트 전략
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s2 = prs.slides.add_slide(blank)
add_rect(s2, 0, 0, 13.33, 7.5, LGRAY)
header_bar(s2, "테스트 전략", "단위 테스트 선정 기준 및 작성 방침")
footer(s2, 2)

# 3개 카드
cards = [
    (BLUE,   "선정 기준",
     ["프로젝트 코드를 직접 검증 (서드파티 라이브러리 테스트 제외)",
      "독립 함수/헤더온리 모듈 우선 — ROS2·raisim 의존성 없이 실행 가능",
      "버그 영향도 큰 알고리즘: A*, 체크섬, 좌표 변환 등 ★★★ 항목"]),
    (GREEN,  "기존 테스트 활용",
     ["test_ublox_gps_init.cpp — 시리얼 초기화 실패·isInitialized() 검증 유지",
      "test_map_services.cpp (fast_lio) — 지도 I/O·다운샘플링 이미 커버",
      "신규 테스트는 기존 미커버 영역(checksum 엣지 케이스, map_graph) 집중"]),
    (ORANGE, "테스트 제외 항목",
     ["서드파티 라이브러리 직접 검증 (GeographicLib, PCL VoxelGrid, Ouster SDK3)",
      "분리 불가 인라인 로직 (navigation_main.cpp subprocess, ImGui draw loop)",
      "하드웨어 의존 통합 테스트 (실 LiDAR, u-Blox, raisim::ArticulatedSystem*)"]),
]
for i, (color, title, items) in enumerate(cards):
    x = 0.35 + i * 4.33
    add_rect(s2, x, 1.55, 4.1, 5.45, WHITE)
    add_rect(s2, x, 1.55, 4.1, 0.38, color)
    add_text_box(s2, title, x+0.15, 1.58, 3.8, 0.33,
                 size=13, bold=True, color=WHITE)
    tb = s2.shapes.add_textbox(Inches(x+0.15), Inches(2.05),
                                Inches(3.8), Inches(4.7))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(11)
        run.font.color.rgb = BLACK

# 빌드 플래그 배너
add_rect(s2, 0.35, 7.0, 12.6, 0.15, NAVY)
add_text_box(s2, "빌드 플래그:  -DRAISIN_BUILD_TEST=ON   |   GTest PRIVATE 링크   |   C++20 std",
             0.5, 7.01, 12, 0.13, size=8, color=WHITE)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 슬라이드 3: 구현된 테스트 목록
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s3 = prs.slides.add_slide(blank)
add_rect(s3, 0, 0, 13.33, 7.5, LGRAY)
header_bar(s3, "구현된 테스트 목록", "신규 작성 테스트 케이스 요약 (기존 2건 + 신규 20건)")
footer(s3, 3)

# 왼쪽: map_graph 테스트
add_rect(s3, 0.3, 1.5, 6.2, 5.6, WHITE)
add_rect(s3, 0.3, 1.5, 6.2, 0.35, BLUE)
add_text_box(s3, "raisin_autonomy_plugin  —  test_map_graph.cpp",
             0.45, 1.53, 5.9, 0.3, size=11, bold=True, color=WHITE)

map_tests = [
    ("ProjectPointOnSegmentClampsToEndpoint",    "기존", "끝점 클램핑"),
    ("FindsInteriorProjectionBeforeNearestNode",  "기존", "엣지 중간 투영 우선"),
    ("AStarReturnsShortestPath",                  "기존", "최단 경로 (우회 비교)"),
    ("RefineWaypointsInjectsGraphNodes",          "기존", "경로 정제 노드 삽입"),
    ("AStarStartEqualsGoalReturnsSingleNode",     "신규", "시작==목표 → 크기 1 경로"),
    ("AStarUnreachableGoalReturnsEmpty",          "신규", "단절 그래프 → 빈 경로"),
    ("ProjectPointOnSegmentDegenerateEdge",       "신규", "길이 0 엣지 처리"),
    ("ProjectPointOnSegmentBeforeStart",          "신규", "세그먼트 앞쪽 점 클램핑"),
    ("FindNearestFallsBackToNodeWhenProjectionNearEndpoint", "신규", "T_MIN 클리핑 → 노드 폴백"),
    ("BuildAdjacencySkipsNegativeCostEdge",       "신규", "음수 비용 엣지 필터"),
    ("BuildAdjacencyMapsMultipleValidEdges",      "신규", "유효 엣지 복수 매핑"),
    ("RefineWaypointsReturnsFalseForGpsFramePairs","신규", "GPS 프레임 쌍 스킵"),
]
tb3l = s3.shapes.add_textbox(Inches(0.45), Inches(1.95), Inches(5.9), Inches(5.1))
tf3l = tb3l.text_frame; tf3l.word_wrap = False
first = True
for name, kind, desc in map_tests:
    p = tf3l.paragraphs[0] if first else tf3l.add_paragraph()
    first = False
    p.space_before = Pt(1)
    run = p.add_run()
    tag_str = "[기존]" if kind == "기존" else "[신규]"
    run.text = f"{tag_str}  {desc}"
    run.font.size = Pt(9.5)
    run.font.color.rgb = DGRAY if kind == "기존" else BLACK
    run.font.bold = (kind == "신규")

# 오른쪽: ublox 테스트
add_rect(s3, 6.83, 1.5, 6.2, 5.6, WHITE)
add_rect(s3, 6.83, 1.5, 6.2, 0.35, GREEN)
add_text_box(s3, "raisin_ublox_plugin  —  test_ublox_utils.cpp",
             6.98, 1.53, 5.9, 0.3, size=11, bold=True, color=WHITE)

ublox_tests = [
    ("UbloxGpsInitFailureTest.*",                    "기존", "시리얼 초기화 실패 / isInitialized()"),
    ("FletcherChecksumMatchesKnownReference",        "수정", "8바이트 참조값 수정 (0x5b→0xAC)"),
    ("FletcherChecksumEmptyPayloadIsZero",           "신규", "빈 페이로드 → ck_a=ck_b=0"),
    ("FletcherChecksumSingleByteEqualsItself",       "신규", "단일 바이트 → ck_a==ck_b==byte"),
    ("FletcherChecksumBitFlipChangesResult",         "신규", "1비트 변조 → 체크섬 불일치"),
    ("FletcherChecksumBIsOrderDependent",            "신규", "바이트 순서 → ck_b 다름"),
    ("FletcherChecksum16BitWrapperPacksBothBytes",   "신규", "uint16_t 래퍼 패킹 검증"),
    ("NmeaChecksumUsesXorAcrossSentenceBody",        "수정", "GPGGA XOR 수정 (69→77)"),
    ("NmeaChecksumEmptyStringIsZero",               "신규", "빈 문자열 → \"0\""),
    ("NmeaChecksumSingleCharacter",                 "신규", "단일 문자 'A' → \"41\""),
    ("NmeaChecksumSameCharTwiceIsZero",             "신규", "동일 문자×2 XOR → \"0\""),
]
tb3r = s3.shapes.add_textbox(Inches(6.98), Inches(1.95), Inches(5.9), Inches(5.1))
tf3r = tb3r.text_frame; tf3r.word_wrap = False
first = True
for name, kind, desc in ublox_tests:
    p = tf3r.paragraphs[0] if first else tf3r.add_paragraph()
    first = False
    p.space_before = Pt(1)
    run = p.add_run()
    if kind == "기존":
        tag_str = "[기존]"; col = DGRAY; bld = False
    elif kind == "수정":
        tag_str = "[수정]"; col = ORANGE; bld = True
    else:
        tag_str = "[신규]"; col = BLACK; bld = True
    run.text = f"{tag_str}  {desc}"
    run.font.size = Pt(9.5)
    run.font.color.rgb = col
    run.font.bold = bld

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 슬라이드 4: 테스트 실행 결과
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s4 = prs.slides.add_slide(blank)
add_rect(s4, 0, 0, 13.33, 7.5, LGRAY)
header_bar(s4, "테스트 실행 결과", "cmake-build-release / Ninja / GTest")
footer(s4, 4)

# 요약 배너
add_rect(s4, 0.3, 1.5, 12.73, 0.72, GREEN)
add_text_box(s4, "전체 24개 테스트  PASSED  (0 FAILED)   —   총 실행 시간 < 1ms",
             0.5, 1.6, 12.3, 0.52, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 패키지별 결과 표
results = [
    # pkg,                      file,                       suite_cnt, test_cnt, pass_cnt, fail_cnt, time
    ("raisin_autonomy_plugin",  "test_map_graph.cpp",       1, 12, 12, 0, "0 ms"),
    ("raisin_ublox_plugin",     "test_ublox_gps_init.cpp",  1,  2,  2, 0, "0 ms"),
    ("raisin_ublox_plugin",     "test_ublox_utils.cpp",     1, 10, 10, 0, "0 ms"),
]

# 표 헤더
hdr_cols = ["패키지", "테스트 파일", "스위트", "전체", "PASS", "FAIL", "시간"]
col_widths = [3.3, 3.5, 0.7, 0.7, 0.7, 0.7, 0.8]
col_starts = [0.3]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

y_hdr = 2.4
add_rect(s4, 0.3, y_hdr, 10.4, 0.33, NAVY)
for j, (col, cw, cx) in enumerate(zip(hdr_cols, col_widths, col_starts)):
    align = PP_ALIGN.CENTER if j >= 2 else PP_ALIGN.LEFT
    add_text_box(s4, col, cx+0.05, y_hdr+0.04, cw-0.1, 0.25,
                 size=10, bold=True, color=WHITE, align=align)

for i, (pkg, fname, sc, tc, pc, fc, tm) in enumerate(results):
    y_row = y_hdr + 0.35 + i*0.4
    bg = WHITE if i%2==0 else LGRAY
    add_rect(s4, 0.3, y_row, 10.4, 0.38, bg)
    row_data = [pkg, fname, str(sc), str(tc), str(pc), str(fc), tm]
    for j, (val, cw, cx) in enumerate(zip(row_data, col_widths, col_starts)):
        align = PP_ALIGN.CENTER if j >= 2 else PP_ALIGN.LEFT
        col = GREEN if (j==4 and fc==0) else (RED if j==5 and fc>0 else BLACK)
        bold = j >= 4
        add_text_box(s4, val, cx+0.05, y_row+0.06, cw-0.1, 0.28,
                     size=10, bold=bold, color=col, align=align)

# 합계 행
y_sum = y_hdr + 0.35 + len(results)*0.4
add_rect(s4, 0.3, y_sum, 10.4, 0.38, NAVY)
sum_vals = ["합계", "", "3", "24", "24", "0", "< 1ms"]
for j, (val, cw, cx) in enumerate(zip(sum_vals, col_widths, col_starts)):
    align = PP_ALIGN.CENTER if j >= 2 else PP_ALIGN.LEFT
    add_text_box(s4, val, cx+0.05, y_sum+0.06, cw-0.1, 0.28,
                 size=10, bold=True, color=WHITE, align=align)

# 실행 커맨드
add_rect(s4, 0.3, 4.5, 12.73, 1.6, RGBColor(0x1C,0x1C,0x1C))
add_text_box(s4, "# 빌드 명령", 0.5, 4.55, 12.4, 0.25,
             size=9, bold=True, color=RGBColor(0x88,0xBB,0xFF))
add_text_box(s4,
    "cd /home/user/raisin_ws/cmake-build-release\n"
    "cmake .   # cmake 캐시 갱신\n"
    "ninja raisin_autonomy_plugin_unittest raisin_ublox_plugin_unittest\n\n"
    "# 실행 명령\n"
    "./src/raisin_plugin/raisin_autonomy_plugin/raisin_autonomy_plugin_unittest --gtest_color=yes\n"
    "./src/raisin_plugin/raisin_ublox_plugin/raisin_ublox_plugin_unittest  --gtest_color=yes",
    0.5, 4.82, 12.4, 1.2,
    size=8.5, bold=False, color=RGBColor(0xBB,0xFF,0xBB))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 슬라이드 5: 발견된 이슈 & 권고사항
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s5 = prs.slides.add_slide(blank)
add_rect(s5, 0, 0, 13.33, 7.5, LGRAY)
header_bar(s5, "발견된 이슈 및 권고사항", "테스트 과정에서 드러난 코드·테스트 결함")
footer(s5, 5)

# 버그 카드
add_rect(s5, 0.3, 1.5, 12.73, 0.3, RED)
add_text_box(s5, "  기존 테스트 파일에서 발견된 버그 2건 (수정 완료)",
             0.3, 1.52, 12, 0.28, size=12, bold=True, color=WHITE)

bugs = [
    ("BUG-01",
     "FletcherChecksumMatchesKnownReference — ck_b 기대값 오류",
     "8바이트 페이로드에서 ck_b = 0xAC 가 정확한 값.\n"
     "기존 코드의 0x5b 는 앞 5바이트만 계산한 잘못된 참조값.",
     "test_ublox_utils.cpp:22   →  EXPECT_EQ(ck_b, 0xAC)  로 수정"),
    ("BUG-02",
     "NmeaChecksumUsesXorAcrossSentenceBody — XOR 결과값 오류",
     "\"GPGGA,123519\" 전체 XOR = 0x77 = \"77\".\n"
     "기존 코드의 \"69\" 는 실제 연산과 불일치.",
     "test_ublox_utils.cpp     →  EXPECT_EQ(calculateChecksum(...), \"77\")  로 수정"),
]
for i, (bid, title, detail, fix) in enumerate(bugs):
    y = 1.92 + i * 1.55
    add_rect(s5, 0.3, y, 12.73, 1.42, WHITE)
    add_rect(s5, 0.3, y, 0.65, 1.42, RED)
    add_text_box(s5, bid, 0.32, y+0.53, 0.62, 0.35, size=9, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s5, title,  1.05, y+0.05, 11.8, 0.3, size=11, bold=True, color=RED)
    add_text_box(s5, detail, 1.05, y+0.4,  11.8, 0.55, size=9.5, color=BLACK)
    add_rect(s5, 1.05, y+0.99, 11.75, 0.3, RGBColor(0xE8,0xF5,0xE9))
    add_text_box(s5, "✔ " + fix, 1.1, y+1.02, 11.65, 0.25, size=9, bold=True, color=GREEN)

# 권고 사항
add_rect(s5, 0.3, 5.08, 12.73, 0.28, BLUE)
add_text_box(s5, "  다음 단계 권고사항",
             0.3, 5.10, 12, 0.24, size=11, bold=True, color=WHITE)

recs = [
    "convertNavPvtToNavSatAndTwist 변환 테스트 추가 (좌표 스케일, fix quality 분기, nanosecond 처리)",
    "raisin_tf: canTransform / getTransform 체인 합성 단위 테스트 추가 (roots differ → Identity 케이스 포함)",
    "VoxelMap::getHeight 테스트 추가 — 경로 계획 컨트롤러가 직접 의존하는 핵심 쿼리 함수",
    "navigation_main.cpp 의 waypoint 진행 로직을 별도 함수로 추출하여 단위 테스트 가능하게 리팩터링",
    "raisin_gui_map1_window: 롱프레스·측정 상태 머신을 ImGui draw loop에서 분리 → 자동화 테스트 가능",
]
tb5 = s5.shapes.add_textbox(Inches(0.4), Inches(5.42), Inches(12.6), Inches(1.7))
tf5 = tb5.text_frame; tf5.word_wrap = True
first = True
for rec in recs:
    p = tf5.paragraphs[0] if first else tf5.add_paragraph()
    first = False; p.space_before = Pt(2)
    run = p.add_run(); run.text = f"→  {rec}"
    run.font.size = Pt(9.5); run.font.color.rgb = BLACK

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
out = "/home/user/raisin_ws/raisin_unit_test_report.pptx"
prs.save(out)
print(f"저장 완료: {out}")
